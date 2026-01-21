"""
Export Utilities for Aadhaar Policy Impact Prediction System
Handles PDF, Excel, CSV, and PowerPoint exports
"""

import pandas as pd
from datetime import datetime
import io
import base64
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.pdfgen import canvas
import matplotlib.pyplot as plt
import matplotlib
matplotlib.use('Agg')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("Warning: python-pptx not installed. PowerPoint export will be disabled.")

class ExportManager:
    """Manages all export functionality"""
    
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
    
    def _setup_custom_styles(self):
        """Setup custom PDF styles"""
        # Title style
        self.title_style = ParagraphStyle(
            'CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#0d47a1'),
            spaceAfter=30,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        # Heading style
        self.heading_style = ParagraphStyle(
            'CustomHeading',
            parent=self.styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#1565c0'),
            spaceAfter=12,
            spaceBefore=12,
            fontName='Helvetica-Bold'
        )
        
        # Body style
        self.body_style = ParagraphStyle(
            'CustomBody',
            parent=self.styles['Normal'],
            fontSize=11,
            textColor=colors.HexColor('#1e3a5f'),
            spaceAfter=12,
            fontName='Helvetica'
        )
    
    def export_to_pdf(self, result_data, policy_name, viz_image_base64=None):
        """
        Generate PDF report
        
        Args:
            result_data: Dictionary with prediction results
            policy_name: Name of the policy
            viz_image_base64: Base64 encoded visualization image
            
        Returns:
            BytesIO object containing PDF
        """
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter,
                               rightMargin=72, leftMargin=72,
                               topMargin=72, bottomMargin=18)
        
        # Container for PDF elements
        elements = []
        
        # Title
        title = Paragraph("Aadhaar Policy Impact Prediction Report", self.title_style)
        elements.append(title)
        elements.append(Spacer(1, 12))
        
        # Policy Information
        policy_info = f"""
        <b>Policy Name:</b> {policy_name}<br/>
        <b>Policy Date:</b> {result_data['summary']['peak_date']}<br/>
        <b>Report Generated:</b> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}<br/>
        """
        elements.append(Paragraph(policy_info, self.body_style))
        elements.append(Spacer(1, 20))
        
        # Executive Summary
        elements.append(Paragraph("Executive Summary", self.heading_style))
        
        summary = result_data['summary']
        summary_data = [
            ['Metric', 'Value'],
            ['Total People Affected', f"{summary['total_people_affected']:,}"],
            ['Total Enrolments', f"{summary['total_enrolments']:,}"],
            ['Total Updates', f"{summary['total_updates']:,}"],
            ['Peak Date', summary['peak_date']],
            ['Peak Volume', f"{summary['peak_volume']:,}"],
            ['Duration', f"{summary['duration_days']} days"]
        ]
        
        summary_table = Table(summary_data, colWidths=[3*inch, 3*inch])
        summary_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1565c0')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f8fafb')),
            ('TEXTCOLOR', (0, 1), (-1, -1), colors.HexColor('#1e3a5f')),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 10),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e1e8ed')),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#fafbfc')])
        ]))
        
        elements.append(summary_table)
        elements.append(Spacer(1, 20))
        
        # Regional Impact Analysis
        elements.append(Paragraph("Regional Impact Analysis", self.heading_style))
        
        regional_data = [['State', 'Enrolments', 'Updates', 'Total Impact', '% Increase', 'Risk']]
        for region in result_data['regional_impact'][:15]:  # Top 15 states
            regional_data.append([
                region['state'],
                f"{region['predicted_enrolments']:,}",
                f"{region['predicted_updates']:,}",
                f"{region['total_impact']:,}",
                f"{region['pct_increase']}%",
                region['risk_level']
            ])
        
        regional_table = Table(regional_data, colWidths=[1.5*inch, 1*inch, 1*inch, 1.2*inch, 0.8*inch, 0.8*inch])
        regional_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1565c0')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 9),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f8fafb')),
            ('TEXTCOLOR', (0, 1), (-1, -1), colors.HexColor('#1e3a5f')),
            ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#e1e8ed')),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#fafbfc')])
        ]))
        
        elements.append(regional_table)
        elements.append(Spacer(1, 20))
        
        # Add visualization if available
        if viz_image_base64:
            elements.append(PageBreak())
            elements.append(Paragraph("Visual Analysis", self.heading_style))
            
            # Decode base64 image
            img_data = base64.b64decode(viz_image_base64.split(',')[1])
            img_buffer = io.BytesIO(img_data)
            
            img = Image(img_buffer, width=6*inch, height=4.5*inch)
            elements.append(img)
        
        # Risk Assessment
        elements.append(PageBreak())
        elements.append(Paragraph("Risk Assessment", self.heading_style))
        
        risk_levels = result_data['risk_assessment']
        risk_text = f"""
        <b>High Risk States ({len(risk_levels['high_risk_states'])}):</b><br/>
        {', '.join(risk_levels['high_risk_states'][:10]) if risk_levels['high_risk_states'] else 'None'}<br/><br/>
        
        <b>Medium Risk States ({len(risk_levels['medium_risk_states'])}):</b><br/>
        {', '.join(risk_levels['medium_risk_states'][:10]) if risk_levels['medium_risk_states'] else 'None'}<br/><br/>
        
        <b>Low Risk States ({len(risk_levels['low_risk_states'])}):</b><br/>
        {', '.join(risk_levels['low_risk_states'][:10]) if risk_levels['low_risk_states'] else 'None'}
        """
        elements.append(Paragraph(risk_text, self.body_style))
        
        # Footer
        elements.append(Spacer(1, 30))
        footer_text = "<i>Generated by Aadhaar Policy Impact Prediction System</i>"
        footer_style = ParagraphStyle('Footer', parent=self.body_style, 
                                     fontSize=9, textColor=colors.grey, alignment=TA_CENTER)
        elements.append(Paragraph(footer_text, footer_style))
        
        # Build PDF
        doc.build(elements)
        buffer.seek(0)
        return buffer
    
    def export_to_excel(self, result_data, policy_name):
        """
        Generate Excel file with multiple sheets
        
        Args:
            result_data: Dictionary with prediction results
            policy_name: Name of the policy
            
        Returns:
            BytesIO object containing Excel file
        """
        buffer = io.BytesIO()
        
        with pd.ExcelWriter(buffer, engine='openpyxl') as writer:
            # Summary Sheet
            summary = result_data['summary']
            summary_df = pd.DataFrame([
                ['Policy Name', policy_name],
                ['Policy Date', summary['peak_date']],
                ['Report Generated', datetime.now().strftime('%Y-%m-%d %H:%M:%S')],
                ['', ''],
                ['Total People Affected', summary['total_people_affected']],
                ['Total Enrolments', summary['total_enrolments']],
                ['Total Updates', summary['total_updates']],
                ['Peak Date', summary['peak_date']],
                ['Peak Volume', summary['peak_volume']],
                ['Duration (days)', summary['duration_days']]
            ], columns=['Metric', 'Value'])
            summary_df.to_excel(writer, sheet_name='Summary', index=False)
            
            # Regional Impact Sheet
            regional_df = pd.DataFrame(result_data['regional_impact'])
            regional_df.to_excel(writer, sheet_name='Regional Impact', index=False)
            
            # Daily Impact Sheet
            daily_df = pd.DataFrame(result_data['daily_impact'])
            daily_df.to_excel(writer, sheet_name='Daily Impact', index=False)
            
            # Risk Assessment Sheet
            risk_data = []
            for level in ['high_risk_states', 'medium_risk_states', 'low_risk_states']:
                for state in result_data['risk_assessment'][level]:
                    risk_data.append({
                        'State': state,
                        'Risk Level': level.replace('_risk_states', '').title()
                    })
            risk_df = pd.DataFrame(risk_data)
            risk_df.to_excel(writer, sheet_name='Risk Assessment', index=False)
            
            # Format worksheets
            for sheet_name in writer.sheets:
                worksheet = writer.sheets[sheet_name]
                for column in worksheet.columns:
                    max_length = 0
                    column = [cell for cell in column]
                    for cell in column:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(cell.value)
                        except:
                            pass
                    adjusted_width = (max_length + 2)
                    worksheet.column_dimensions[column[0].column_letter].width = adjusted_width
        
        buffer.seek(0)
        return buffer
    
    def export_to_csv(self, result_data):
        """
        Generate CSV file with regional impact data
        
        Args:
            result_data: Dictionary with prediction results
            
        Returns:
            BytesIO object containing CSV file
        """
        buffer = io.BytesIO()
        
        # Create DataFrame from regional impact
        df = pd.DataFrame(result_data['regional_impact'])
        
        # Write to CSV
        df.to_csv(buffer, index=False, encoding='utf-8')
        buffer.seek(0)
        return buffer
    
    def export_to_powerpoint(self, result_data, policy_name, viz_image_base64=None):
        """
        Generate PowerPoint presentation
        
        Args:
            result_data: Dictionary with prediction results
            policy_name: Name of the policy
            viz_image_base64: Base64 encoded visualization image
            
        Returns:
            BytesIO object containing PowerPoint file
        """
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not installed. Install it with: pip install python-pptx")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define colors
        blue_dark = RGBColor(13, 71, 161)
        blue_light = RGBColor(21, 101, 192)
        
        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Aadhaar Policy Impact Prediction"
        subtitle.text = f"{policy_name}\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Format title
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = blue_dark
        
        # Slide 2: Executive Summary
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Executive Summary"
        title_shape.text_frame.paragraphs[0].font.color.rgb = blue_dark
        
        tf = body_shape.text_frame
        summary = result_data['summary']
        
        p = tf.paragraphs[0]
        p.text = f"Total People Affected: {summary['total_people_affected']:,}"
        p.level = 0
        
        for text in [
            f"Enrolments: {summary['total_enrolments']:,}",
            f"Updates: {summary['total_updates']:,}",
            f"Peak Date: {summary['peak_date']}",
            f"Peak Volume: {summary['peak_volume']:,}",
            f"Duration: {summary['duration_days']} days"
        ]:
            p = tf.add_paragraph()
            p.text = text
            p.level = 0
            p.font.size = Pt(18)
        
        # Slide 3: Regional Impact
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        title_shape.text = "Top 10 Affected States"
        title_shape.text_frame.paragraphs[0].font.color.rgb = blue_dark
        
        # Add table
        rows = 11
        cols = 4
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        
        table = shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(1.8)
        table.columns[2].width = Inches(1.8)
        table.columns[3].width = Inches(1.4)
        
        # Header row
        headers = ['State', 'Enrolments', 'Updates', 'Risk Level']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = blue_light
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        # Data rows
        for i, region in enumerate(result_data['regional_impact'][:10], start=1):
            table.cell(i, 0).text = region['state']
            table.cell(i, 1).text = f"{region['predicted_enrolments']:,}"
            table.cell(i, 2).text = f"{region['predicted_updates']:,}"
            table.cell(i, 3).text = region['risk_level']
            
            for j in range(4):
                table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(10)
        
        # Slide 4: Visualization
        if viz_image_base64:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
            title_frame = title_box.text_frame
            title_frame.text = "Visual Analysis"
            title_frame.paragraphs[0].font.size = Pt(32)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = blue_dark
            
            # Decode and add image
            img_data = base64.b64decode(viz_image_base64.split(',')[1])
            img_buffer = io.BytesIO(img_data)
            
            left = Inches(1)
            top = Inches(1.5)
            pic = slide.shapes.add_picture(img_buffer, left, top, width=Inches(8))
        
        # Slide 5: Risk Assessment
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = "Risk Assessment"
        title_shape.text_frame.paragraphs[0].font.color.rgb = blue_dark
        
        tf = body_shape.text_frame
        risk_levels = result_data['risk_assessment']
        
        p = tf.paragraphs[0]
        p.text = f"High Risk States ({len(risk_levels['high_risk_states'])})"
        p.font.bold = True
        p.font.color.rgb = RGBColor(198, 40, 40)
        
        if risk_levels['high_risk_states']:
            p = tf.add_paragraph()
            p.text = ', '.join(risk_levels['high_risk_states'][:5])
            p.level = 1
            p.font.size = Pt(14)
        
        p = tf.add_paragraph()
        p.text = f"Medium Risk States ({len(risk_levels['medium_risk_states'])})"
        p.font.bold = True
        p.font.color.rgb = RGBColor(239, 108, 0)
        
        if risk_levels['medium_risk_states']:
            p = tf.add_paragraph()
            p.text = ', '.join(risk_levels['medium_risk_states'][:5])
            p.level = 1
            p.font.size = Pt(14)
        
        # Save to buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer

# Singleton instance
export_manager = ExportManager()
